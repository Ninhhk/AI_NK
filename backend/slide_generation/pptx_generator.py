import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import datetime
import logging
import json

logger = logging.getLogger(__name__)

class PowerPointGenerator:
    def __init__(self):
        self.prs = Presentation()

    def generate_presentation(self, slides_data: list, output_path: Path) -> str:
        """Generate a PowerPoint presentation from the slides data.
        
        Args:
            slides_data: List of slide dictionaries with title_text and content fields
            output_path: Path where to save the presentation
            
        Returns:
            Path to the generated presentation
        """
        try:
            # Create a new Presentation object for each generation
            self.prs = Presentation()
            
            # Process all slides - including the first one
            for slide_data in slides_data:
                self.add_slide(slide_data)

            # Ensure the output path exists
            output_path.parent.mkdir(parents=True, exist_ok=True)
            
            # Save the presentation with proper encoding
            self.prs.save(str(output_path))
            logger.info(f"Successfully saved presentation to {output_path}")
            return str(output_path)

        except Exception as e:
            logger.error(f"Error creating PowerPoint presentation: {str(e)}")
            raise

    def add_slide(self, slide_data: dict):
        """Add a slide to the presentation with title and bullet points."""
        try:
            # Use a slide layout with title and content
            slide_layout = self.prs.slide_layouts[1]  # Layout 1 is title and content
            slide = self.prs.slides.add_slide(slide_layout)
            
            # Add title
            if "title_text" in slide_data:
                title = slide.shapes.title
                title.text = slide_data["title_text"]
            
            # Check if content is available (as list)
            if "content" in slide_data and isinstance(slide_data["content"], list):
                # Get the content placeholder (usually index 1)
                content_placeholder = None
                for shape in slide.placeholders:
                    if hasattr(shape, 'placeholder_format') and shape.placeholder_format.idx == 1:  # 1 is the content placeholder
                        content_placeholder = shape
                        break
                
                # If no content placeholder found, try to get by index
                if not content_placeholder:
                    try:
                        content_placeholder = slide.placeholders[1]
                    except:
                        logger.warning("Couldn't find content placeholder, using text box")
                        # Create a text box if placeholder not found
                        left = Inches(1)
                        top = Inches(2)
                        width = Inches(8)
                        height = Inches(4)
                        content_placeholder = slide.shapes.add_textbox(left, top, width, height)
                
                # Get the text frame to add bullet points
                text_frame = content_placeholder.text_frame
                text_frame.clear()  # Clear any default content
                
                # Add each content item as a bullet point
                for i, item in enumerate(slide_data["content"]):
                    if i == 0:
                        # Use the first paragraph that already exists
                        paragraph = text_frame.paragraphs[0]
                    else:
                        # Add additional paragraphs for remaining items
                        paragraph = text_frame.add_paragraph()
                    
                    # Set bullet style
                    paragraph.level = 0  # Top level bullet
                    
                    # Clean the item text (remove leading bullet if present)
                    text = item.strip()
                    if text.startswith("-") or text.startswith("•") or text.startswith("*"):
                        text = text[1:].strip()
                    
                    paragraph.text = text
                    
                    # Set font properties for better appearance
                    for run in paragraph.runs:
                        run.font.size = Pt(18)
            
            # Handle text field if provided instead of content
            elif "text" in slide_data and slide_data["text"]:
                # Process the same way as above but with the text split by lines
                try:
                    content_placeholder = slide.placeholders[1]
                    text_frame = content_placeholder.text_frame
                    text_frame.clear()
                    
                    lines = [line.strip() for line in slide_data["text"].split("\n") if line.strip()]
                    
                    for i, line in enumerate(lines):
                        if i == 0:
                            paragraph = text_frame.paragraphs[0]
                        else:
                            paragraph = text_frame.add_paragraph()
                        
                        paragraph.level = 0
                        
                        # Clean the line
                        if line.startswith("-") or line.startswith("•") or line.startswith("*"):
                            line = line[1:].strip()
                        
                        paragraph.text = line
                        
                        for run in paragraph.runs:
                            run.font.size = Pt(18)
                except Exception as e:
                    logger.error(f"Error processing text field: {e}")
            
            # Add images if present
            if "images" in slide_data:
                for img_path in slide_data.get("images", []):
                    if os.path.exists(img_path):
                        left = Inches(1)
                        top = Inches(3)
                        width = Inches(4)
                        slide.shapes.add_picture(img_path, left, top, width=width)
                
        except Exception as e:
            logger.error(f"Error adding slide: {str(e)}")
            raise Exception(f"Error adding slide: {str(e)}")
